"""
PPTX document loader with intelligent processing.
"""

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table

logger = logging.getLogger(__name__)


class PPTXLoader:
    """Loads and intelligently processes content from PPTX files."""

    def __init__(self):
        """Initialize the PPTX loader."""
        self.max_slide_content_length = 2000  # Limit per slide for very verbose slides

    def _extract_text_from_shape(self, shape) -> str:
        """
        Extract text from a shape, handling different shape types.

        Args:
            shape: PowerPoint shape object

        Returns:
            Extracted text content
        """
        try:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip()
            elif hasattr(shape, "text_frame") and shape.text_frame:
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ""
                    for run in paragraph.runs:
                        if run.text:
                            para_text += run.text
                    if para_text.strip():
                        paragraphs.append(para_text.strip())
                return "\n".join(paragraphs)
        except Exception as e:
            logger.debug(f"Could not extract text from shape: {e}")

        return ""

    def _extract_table_content(self, table: Table) -> str:
        """
        Extract content from a table shape.

        Args:
            table: PowerPoint table object

        Returns:
            Formatted table content
        """
        try:
            table_content = []
            table_content.append("--- Table Content ---")

            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    # Clean up cell text
                    cell_text = " ".join(cell_text.split())  # Remove extra whitespace
                    row_data.append(cell_text if cell_text else "[empty]")

                # Format as table row
                if row_idx == 0:  # Assume first row is header
                    table_content.append("HEADERS: " + " | ".join(row_data))
                else:
                    table_content.append(f"Row {row_idx}: " + " | ".join(row_data))

            return "\n".join(table_content)

        except Exception as e:
            logger.debug(f"Could not extract table content: {e}")
            return ""

    def _analyze_slide_layout(self, slide) -> Dict[str, Any]:
        """
        Analyze slide layout and structure.

        Args:
            slide: PowerPoint slide object

        Returns:
            Dictionary with layout analysis
        """
        analysis = {
            "has_title": False,
            "has_content": False,
            "has_images": False,
            "has_tables": False,
            "has_charts": False,
            "shape_count": 0,
            "text_shapes": 0,
            "layout_name": "Unknown",
        }

        try:
            # Get layout name if available
            if hasattr(slide, "slide_layout") and hasattr(slide.slide_layout, "name"):
                analysis["layout_name"] = slide.slide_layout.name
        except Exception:
            pass

        analysis["shape_count"] = len(slide.shapes)

        for shape in slide.shapes:
            try:
                # Check for title
                if hasattr(shape, "name") and "title" in shape.name.lower():
                    analysis["has_title"] = True

                # Check shape type
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    analysis["has_images"] = True
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    analysis["has_tables"] = True
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    analysis["has_charts"] = True

                # Check for text content
                if hasattr(shape, "text") or hasattr(shape, "text_frame"):
                    text = self._extract_text_from_shape(shape)
                    if text:
                        analysis["text_shapes"] += 1
                        analysis["has_content"] = True

            except Exception as e:
                logger.debug(f"Error analyzing shape: {e}")
                continue

        return analysis

    def _process_slide(self, slide, slide_number: int) -> str:
        """
        Process a single slide and extract all content.

        Args:
            slide: PowerPoint slide object
            slide_number: Slide number (1-based)

        Returns:
            Formatted slide content
        """
        slide_content = []
        slide_content.append(f"=== SLIDE {slide_number} ===")

        # Analyze slide layout
        layout_analysis = self._analyze_slide_layout(slide)
        slide_content.append(f"Layout: {layout_analysis['layout_name']}")

        # Add layout context
        layout_info = []
        if layout_analysis["has_title"]:
            layout_info.append("title")
        if layout_analysis["has_content"]:
            layout_info.append("content")
        if layout_analysis["has_images"]:
            layout_info.append("images")
        if layout_analysis["has_tables"]:
            layout_info.append("tables")
        if layout_analysis["has_charts"]:
            layout_info.append("charts")

        if layout_info:
            slide_content.append(f"Contains: {', '.join(layout_info)}")

        slide_content.append("")  # Empty line

        # Extract content from all shapes
        title_found = False
        content_parts = []

        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # Handle different shape types
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_content = self._extract_table_content(shape.table)
                    if table_content:
                        content_parts.append(table_content)

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Note presence of image
                    content_parts.append(
                        f"[Image: {getattr(shape, 'name', f'Image_{shape_idx}')}]"
                    )

                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    # Note presence of chart
                    content_parts.append(
                        f"[Chart: {getattr(shape, 'name', f'Chart_{shape_idx}')}]"
                    )

                else:
                    # Extract text content
                    text_content = self._extract_text_from_shape(shape)
                    if text_content:
                        # Try to identify if this is a title
                        is_title = (
                            not title_found
                            and (
                                hasattr(shape, "name") and "title" in shape.name.lower()
                            )
                            or len(text_content) < 100  # Short text likely to be title
                        )

                        if is_title and not title_found:
                            slide_content.append(f"TITLE: {text_content}")
                            title_found = True
                        else:
                            content_parts.append(text_content)

            except Exception as e:
                logger.debug(
                    f"Error processing shape {shape_idx} on slide {slide_number}: {e}"
                )
                continue

        # Add all content parts
        if content_parts:
            slide_content.append("CONTENT:")
            slide_content.extend(content_parts)

        # Limit content length per slide
        full_slide_content = "\n".join(slide_content)
        if len(full_slide_content) > self.max_slide_content_length:
            full_slide_content = (
                full_slide_content[: self.max_slide_content_length]
                + "\n[Content truncated...]"
            )

        return full_slide_content

    def _extract_presentation_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """
        Extract metadata from the presentation.

        Args:
            prs: PowerPoint presentation object

        Returns:
            Dictionary with presentation metadata
        """
        metadata = {
            "slide_count": len(prs.slides),
            "has_master_slides": False,
            "unique_layouts": set(),
            "total_shapes": 0,
            "slides_with_content": 0,
            "slides_with_images": 0,
            "slides_with_tables": 0,
            "slides_with_charts": 0,
        }

        try:
            # Analyze all slides for statistics
            for slide in prs.slides:
                analysis = self._analyze_slide_layout(slide)
                metadata["total_shapes"] += analysis["shape_count"]
                metadata["unique_layouts"].add(analysis["layout_name"])

                if analysis["has_content"]:
                    metadata["slides_with_content"] += 1
                if analysis["has_images"]:
                    metadata["slides_with_images"] += 1
                if analysis["has_tables"]:
                    metadata["slides_with_tables"] += 1
                if analysis["has_charts"]:
                    metadata["slides_with_charts"] += 1

            # Convert set to count
            metadata["unique_layout_count"] = len(metadata["unique_layouts"])
            metadata["unique_layouts"] = list(metadata["unique_layouts"])

        except Exception as e:
            logger.debug(f"Error extracting presentation metadata: {e}")

        return metadata

    def _generate_presentation_summary(
        self, metadata: Dict[str, Any], filename: str
    ) -> str:
        """
        Generate a comprehensive summary of the presentation.

        Args:
            metadata: Presentation metadata
            filename: Original filename

        Returns:
            Formatted presentation summary
        """
        summary = []
        summary.append(f"=== PRESENTATION ANALYSIS: {filename} ===")
        summary.append(f"Total slides: {metadata['slide_count']}")
        summary.append(f"Total shapes/elements: {metadata['total_shapes']}")
        summary.append(f"Unique layouts used: {metadata['unique_layout_count']}")

        if metadata["unique_layouts"]:
            summary.append(f"Layout types: {', '.join(metadata['unique_layouts'])}")

        # Content distribution
        summary.append("")
        summary.append("Content Distribution:")
        summary.append(
            f"  - Slides with text content: {metadata['slides_with_content']}"
        )
        summary.append(f"  - Slides with images: {metadata['slides_with_images']}")
        summary.append(f"  - Slides with tables: {metadata['slides_with_tables']}")
        summary.append(f"  - Slides with charts: {metadata['slides_with_charts']}")

        # Presentation insights
        summary.append("")
        summary.append("Presentation Insights:")

        if metadata["slides_with_images"] > metadata["slide_count"] * 0.5:
            summary.append("  - Image-heavy presentation (visual/design focus)")

        if metadata["slides_with_tables"] > 0:
            summary.append("  - Contains structured data in tables")

        if metadata["slides_with_charts"] > 0:
            summary.append("  - Contains data visualizations/charts")

        content_ratio = (
            metadata["slides_with_content"] / metadata["slide_count"]
            if metadata["slide_count"] > 0
            else 0
        )
        if content_ratio > 0.8:
            summary.append("  - Content-rich presentation (text-focused)")
        elif content_ratio < 0.3:
            summary.append("  - Minimal text content (likely visual/template slides)")

        return "\n".join(summary)

    def load(self, file_path: Path) -> Optional[str]:
        """
        Load and intelligently process content from a PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Processed text content or None if failed
        """
        try:
            # Load presentation
            prs = Presentation(str(file_path))

            if len(prs.slides) == 0:
                logger.warning(f"No slides found in PPTX: {file_path}")
                return None

            # Extract metadata and create summary
            metadata = self._extract_presentation_metadata(prs)
            presentation_summary = self._generate_presentation_summary(
                metadata, file_path.name
            )

            # Process all slides
            slide_contents = []
            slide_contents.append(presentation_summary)
            slide_contents.append("")  # Separator

            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = self._process_slide(slide, slide_idx)
                slide_contents.append(slide_content)
                slide_contents.append("")  # Separator between slides

            # Add final summary for context
            slide_contents.append(
                self._generate_content_insights(slide_contents, metadata)
            )

            full_content = "\n".join(slide_contents)

            logger.info(
                f"Successfully loaded PPTX: {file_path} "
                f"({metadata['slide_count']} slides, {metadata['total_shapes']} elements)"
            )

            return full_content

        except Exception as e:
            logger.error(f"Failed to load PPTX {file_path}: {e}")
            return None

    def _generate_content_insights(
        self, slide_contents: List[str], metadata: Dict[str, Any]
    ) -> str:
        """
        Generate insights about the presentation content.

        Args:
            slide_contents: List of slide content strings
            metadata: Presentation metadata

        Returns:
            Content insights text
        """
        insights = []
        insights.append("=== CONTENT INSIGHTS ===")

        # Analyze text patterns
        all_text = " ".join(slide_contents).lower()

        # Detect presentation type
        business_keywords = [
            "revenue",
            "profit",
            "growth",
            "market",
            "sales",
            "customer",
            "strategy",
            "roi",
            "kpi",
        ]
        technical_keywords = [
            "algorithm",
            "data",
            "system",
            "architecture",
            "implementation",
            "api",
            "database",
        ]
        educational_keywords = [
            "learn",
            "understand",
            "example",
            "definition",
            "concept",
            "theory",
            "practice",
        ]

        business_score = sum(1 for keyword in business_keywords if keyword in all_text)
        technical_score = sum(
            1 for keyword in technical_keywords if keyword in all_text
        )
        educational_score = sum(
            1 for keyword in educational_keywords if keyword in all_text
        )

        presentation_type = "General"
        if business_score > technical_score and business_score > educational_score:
            presentation_type = "Business/Commercial"
        elif technical_score > business_score and technical_score > educational_score:
            presentation_type = "Technical/Engineering"
        elif educational_score > business_score and educational_score > technical_score:
            presentation_type = "Educational/Training"

        insights.append(f"Presentation type: {presentation_type}")

        # Content complexity
        avg_shapes_per_slide = (
            metadata["total_shapes"] / metadata["slide_count"]
            if metadata["slide_count"] > 0
            else 0
        )
        if avg_shapes_per_slide > 10:
            insights.append("Content complexity: High (dense slides)")
        elif avg_shapes_per_slide > 5:
            insights.append("Content complexity: Medium")
        else:
            insights.append("Content complexity: Low (simple slides)")

        # Structural insights
        if metadata["slides_with_tables"] > 0 and metadata["slides_with_charts"] > 0:
            insights.append(
                "Structure: Data-driven presentation with both tables and visualizations"
            )
        elif metadata["slides_with_images"] > metadata["slide_count"] * 0.7:
            insights.append("Structure: Visual-heavy presentation")
        elif metadata["slides_with_content"] == metadata["slide_count"]:
            insights.append("Structure: Text-focused presentation")

        return "\n".join(insights)
